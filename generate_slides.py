#!/usr/bin/env python3
"""Generate PyCon Namibia talk slides as .pptx"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Spotify brand-ish colors
BLACK = RGBColor(0x19, 0x19, 0x19)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x1D, 0xB9, 0x54)
GRAY = RGBColor(0xB3, 0xB3, 0xB3)
DARK_GRAY = RGBColor(0x53, 0x53, 0x53)


def set_slide_bg(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_slide(prs, title, subtitle, notes=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, BLACK)

    # Title
    left = Inches(1)
    top = Inches(2.0)
    width = Inches(8)
    height = Inches(1.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = GREEN
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    top2 = Inches(3.6)
    txBox2 = slide.shapes.add_textbox(left, top2, width, Inches(1.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(22)
    p2.font.color.rgb = GRAY
    p2.alignment = PP_ALIGN.CENTER

    if notes:
        slide.notes_slide.notes_text_frame.text = notes

    return slide


def add_content_slide(prs, title, bullets, notes=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, BLACK)

    # Title
    left = Inches(0.8)
    top = Inches(0.5)
    width = Inches(8.4)
    txBox = slide.shapes.add_textbox(left, top, width, Inches(1.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = GREEN

    # Bullets
    top2 = Inches(1.7)
    txBox2 = slide.shapes.add_textbox(left, top2, width, Inches(5.0))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
            p.space_before = Pt(12)
        p.text = bullet
        p.font.size = Pt(22)
        p.font.color.rgb = WHITE
        p.level = 0

    if notes:
        slide.notes_slide.notes_text_frame.text = notes

    return slide


def add_section_slide(prs, title, subtitle="", notes=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, BLACK)

    top = Inches(2.5)
    left = Inches(1)
    width = Inches(8)
    txBox = slide.shapes.add_textbox(left, top, width, Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = GREEN
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        top2 = Inches(3.8)
        txBox2 = slide.shapes.add_textbox(left, top2, width, Inches(1.0))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(22)
        p2.font.color.rgb = GRAY
        p2.alignment = PP_ALIGN.CENTER

    if notes:
        slide.notes_slide.notes_text_frame.text = notes

    return slide


def add_end_slide(prs, title, lines, notes=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, BLACK)

    left = Inches(1)
    top = Inches(1.5)
    width = Inches(8)
    txBox = slide.shapes.add_textbox(left, top, width, Inches(1.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = GREEN
    p.alignment = PP_ALIGN.CENTER

    top2 = Inches(3.0)
    txBox2 = slide.shapes.add_textbox(left, top2, width, Inches(3.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
            p.space_before = Pt(8)
        p.text = line
        p.font.size = Pt(20)
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.CENTER

    if notes:
        slide.notes_slide.notes_text_frame.text = notes

    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9

    # ── Slide 1: Title ──
    add_title_slide(
        prs,
        "From Script Kiddie to Spotify",
        "A 30-Year Journey\n\nGijs Molenaar\nPyCon Namibia 2026",
        notes=(
            "Welcome everyone. I'm Gijs Molenaar, a software engineer at Spotify. "
            "Today I want to share my 30-year journey in tech — from a curious teenager "
            "in the Netherlands to working at one of the world's biggest music platforms. "
            "This is an honest story about persistence, curiosity, and a lot of luck."
        ),
    )

    # ── Slide 2: The Record Store ──
    add_content_slide(
        prs,
        "The Record Store",
        [
            "My parents owned a record store in the Netherlands",
            "Music was everywhere growing up",
            "Vinyl, CDs, cassettes — the physical music era",
        ],
        notes=(
            "Let me start at the beginning. My parents had a record store. "
            "Music was the soundtrack of my childhood — literally. Stacks of vinyl, "
            "walls of CDs. This is important because the story comes full circle later. "
            "Keep the record store in mind."
        ),
    )

    # ── Slide 3: The Curious Kid ──
    add_content_slide(
        prs,
        "The Curious Kid",
        [
            "Age 14, volunteering at a local TV station",
            "Discovered the Windows NT 'ping of death'",
            "Just a kid playing with scripts I barely understood",
        ],
        notes=(
            "At 14, I was volunteering at a local TV station. I discovered that their "
            "ISP's Windows NT systems didn't survive a 'ping of death.' I was just a "
            "curious kid playing with security tools — scripts I found online and barely "
            "understood. The classic script kiddie."
        ),
    )

    # ── Slide 4: First Job at 15 ──
    add_content_slide(
        prs,
        "First Job at 15",
        [
            "The ISP heard about 'the kid who knows security'",
            "Instead of calling the police — they called me",
            "Security consultant at 15 years old",
            "⚠️ Don't try this today. I got lucky.",
        ],
        notes=(
            "By complete coincidence, the ISP heard there was a young person who knew "
            "something about computer security. Instead of calling the police, they "
            "called me for help. At 15, I had my first job as a security consultant. "
            "A word of warning: I got lucky. In the early 2000s, law enforcement wasn't "
            "equipped to handle these situations. Today, the same curiosity could land "
            "you in jail. Seriously — don't do what I did."
        ),
    )

    # ── Slide 5: Finding My Path ──
    add_content_slide(
        prs,
        "Finding My Path",
        [
            "MBO in electronics — discovered I didn't want to repair fridges",
            "Got my Cisco CCNA certification",
            "Terrible student → motivated once I found computing",
            "Not knowing what you want is normal",
        ],
        notes=(
            "Around the same time, I enrolled in MBO to study electronics. People had "
            "recommended it — I had no idea what I wanted. I got my CCNA during the "
            "program, but I also realized: I don't want to repair fridges. I don't want "
            "high-voltage equipment. I wanted software. I wanted science. "
            "Here's the thing: I was a terrible student. My marks were bad. But once I "
            "found computing — something I actually cared about — everything changed. "
            "Between 15 and 18, I went from one of the worst performers to someone who "
            "genuinely wanted to excel."
        ),
    )

    # ── Slide 6: Software Engineering ──
    add_content_slide(
        prs,
        "Software Engineering",
        [
            "Bachelor's at Hogeschool van Amsterdam",
            "Bridging the gap from vocational to higher education",
            "Teachers helped me fill the math gaps",
        ],
        notes=(
            "I did my Bachelor's in Software Engineering at the Hogeschool van Amsterdam. "
            "Going from MBO to HBO was a real transition — especially the math. "
            "Teachers at the Hogeschool helped me bridge that gap. This is another example "
            "of people investing in you — if you show willingness to learn, people help."
        ),
    )

    # ── Slide 7: Python: Love at First Sight ──
    add_content_slide(
        prs,
        "Python: Love at First Sight",
        [
            "Using Python since that first ISP job — nearly 30 years",
            "Named my company 'Pythonic'",
            "Prototype in Python, optimize later",
            "Python 2→3 transition, wheels, uv, ruff — I've seen it all",
        ],
        notes=(
            "I've been using Python since my first job at that ISP. It's been my "
            "30-year companion. I named my company Pythonic! I got lucky that Python "
            "became so popular. I've used it for everything: websites, scientific "
            "computing, art projects, machine learning. That pattern — prototype in "
            "Python, optimize later — has served me ever since. I've witnessed the "
            "Python 2 to 3 transition, which was a nightmare in science. And now "
            "the modern tooling renaissance with uv and ruff."
        ),
    )

    # ── Slide 8: Master's in AI ──
    add_content_slide(
        prs,
        "Master's in Artificial Intelligence",
        [
            "University of Amsterdam",
            "Sonic Gesture: hand movements → MIDI music",
            "Prototyped in Python, rewrote in C++ for real-time",
            "[Demo clip placeholder]",
        ],
        notes=(
            "For my Master's thesis at the University of Amsterdam, I built 'Sonic "
            "Gesture' — a computer vision system that translated hand movements into "
            "MIDI music. I prototyped it in Python, then rewrote it in C++ when I needed "
            "real-time video performance. Fun fact: one of my professors was the creator "
            "of MonetDB, a database system I later worked on as a contractor, and which "
            "inspired DuckDB."
        ),
    )

    # ── Slide 9: Open Source ──
    add_content_slide(
        prs,
        "The Power of Open Source",
        [
            "First contributions to Debian — before GitHub existed",
            "Fix a bug, improve docs, package something",
            "Taught me more than any course",
            "Built a portfolio that opened doors",
            "The open source community remembers contributors",
        ],
        notes=(
            "Python and open source were my great equalizers. My first contributions "
            "were to Debian — packaging software, submitting bug reports. This was before "
            "GitHub even existed. Contributing to projects taught me more than any course, "
            "connected me with maintainers worldwide, and built a portfolio that opened "
            "doors. You don't need to write a major feature. Start small."
        ),
    )

    # ── Slide 10: The Power of Mentors ──
    add_content_slide(
        prs,
        "The Power of Mentors",
        [
            "Girlfriend's father: science fiction, Rubik's Cubes, Gödel-Escher-Bach",
            "His brother at SARA → internship → scientific computing",
            "Teachers who helped bridge the math gap",
            "Each person led to the next — you can't plan this",
        ],
        notes=(
            "I got lucky with people. My first mentor was the father of my first "
            "girlfriend — a geeky chemist who introduced me to science fiction, "
            "Rubik's Cubes, and Gödel-Escher-Bach. My parents didn't know this world "
            "existed; he opened the door. His brother worked at SARA, the Dutch national "
            "supercomputer facility, and arranged my internship there. That introduced me "
            "to scientific computing. From there I met people at SURFnet who later became "
            "clients. Each person led to the next. You can't plan this. But you can put "
            "yourself in positions where it might happen."
        ),
    )

    # ── Slide 11: The Record Store Dies ──
    add_content_slide(
        prs,
        "The Record Store Dies",
        [
            "MP3s and online piracy killed the record store",
            "The music industry in crisis",
            "My parents lost their livelihood",
            "Remember this — the story comes back",
        ],
        notes=(
            "Remember the record store I mentioned? MP3s and online copying killed it. "
            "Napster, Limewire, torrents — the music industry was in freefall. My parents "
            "lost their business. This was personal. Keep this in mind, because the story "
            "comes full circle."
        ),
    )

    # ── Slide 12: PhD in South Africa ──
    add_content_slide(
        prs,
        "PhD in South Africa",
        [
            "Rhodes University, Makhanda",
            "The Square Kilometre Array — the world's largest radio telescope",
            "Radio astronomy meets computer science",
            "Working on problems at an incredible scale",
        ],
        notes=(
            "I moved to South Africa for a PhD at Rhodes University, working on the "
            "Square Kilometre Array — one of the most ambitious scientific projects on "
            "Earth. The world's largest radio telescope. Radio astronomy combined with "
            "computer science. The data challenges are enormous — we're talking petabytes "
            "of data per second."
        ),
    )

    # ── Slide 13: Meeting My Wife ──
    add_content_slide(
        prs,
        "Meeting My Wife",
        [
            "Came to South Africa for science",
            "Found love — met a Namibian",
            "Now making Namibia home",
            "She works at Hyphen on green hydrogen",
        ],
        notes=(
            "I came to South Africa for science, and found something unexpected: love. "
            "I met my wife, a Namibian. We've been coming to Namibia regularly ever since, "
            "spending every Christmas with family. Now we're making it official — my wife "
            "has taken a role at Hyphen, working on green hydrogen, and we're moving here "
            "with our child. Namibia is becoming home."
        ),
    )

    # ── Slide 14: The Freelance Years ──
    add_content_slide(
        prs,
        "The Freelance Years",
        [
            "10+ years consulting: ASTRON, SARAO, SETI, SURFnet",
            "Radio astronomy, ML for energy grids, signal processing",
            "Chasing money made me unhappy",
            "Following interesting made me fulfilled",
        ],
        notes=(
            "I spent over 10 years as a freelance consultant, working for organizations "
            "like ASTRON, SARAO, SETI, and SURFnet. Radio astronomy, machine learning "
            "for electricity grids, signal processing. During my consulting years, I made "
            "the mistake of chasing money. It made me unhappy. The best career decisions "
            "I made were about curiosity, not income."
        ),
    )

    # ── Slide 15: Balancing PhD + Freelancing ──
    add_content_slide(
        prs,
        "PhD + Freelancing",
        [
            "4 days client work + 1-2 days PhD = 6-day weeks",
            "Exhausting but worth it",
            "Finished just before COVID hit",
            "Be realistic about the cost of part-time education",
        ],
        notes=(
            "For about two years, I worked four days a week for clients and one or two "
            "days on my PhD. That's a six-day workweek. It was exhausting. I finished "
            "just before COVID hit — I'm not sure I could have kept that pace through a "
            "pandemic. My advice: if you're considering part-time education while working, "
            "be realistic about the cost. It's possible, but it will consume your life "
            "for a while. Make sure the goal is worth it."
        ),
    )

    # ── Slide 16: The Teenage Dream ──
    add_content_slide(
        prs,
        "The Teenage Dream",
        [
            "Spotify launched in the Netherlands in 2010",
            '"That\'s my dream job"',
            "Twelve years later, I got it",
        ],
        notes=(
            "When Spotify launched in the Netherlands in 2010, I thought: that's my dream "
            "job. I'm a music lover, a Python developer, and here was a company that was "
            "changing how the world listens to music. Twelve years later, I got it. "
            "I'm now living my teenage dream — and it pays well too. But the money came "
            "from following interesting work, not the other way around."
        ),
    )

    # ── Slide 17: Full Circle ──
    add_content_slide(
        prs,
        "Full Circle: Music Saved",
        [
            "Parents' record store was killed by MP3s",
            "Spotify saved music — a new revenue model for artists",
            "Now I work there",
            "Life has a sense of humor",
        ],
        notes=(
            "Remember the record store? MP3s killed it. But Spotify saved music by "
            "creating a revenue model for artists. Streaming replaced piracy. And now "
            "I work at the company that made it happen. My parents lost their livelihood "
            "to digital disruption, and I ended up at the company that rebuilt the "
            "industry. Life has a sense of humor."
        ),
    )

    # ── Slide 18: AI Changed Everything ──
    add_content_slide(
        prs,
        "AI Changed Everything",
        [
            "I no longer type code — I talk to Claude with voice",
            "This very talk was created with AI",
            "Orders of magnitude faster",
            "Solving problems I never thought I could tackle",
        ],
        notes=(
            "The era of manual programming is changing. At Spotify, I no longer type "
            "code — I talk to Claude using voice recognition. This very talk, these very "
            "slides — created entirely through conversation with AI. These tools have "
            "changed everything. I can solve problems I never thought I could tackle, "
            "work orders of magnitude faster, and build things that would have taken "
            "weeks in days."
        ),
    )

    # ── Slide 19: What I Built with AI ──
    add_content_slide(
        prs,
        "What I Built with AI",
        [
            "python-snap7 — industrial PLC communication",
            "ableton-claude — control Ableton Live with AI",
            "spectrageist — real-time audio feature extraction",
            "geistwave — experimental music LLM",
            "charcoaloptimizer — AI-powered 3D packing optimization",
            "This CV and talk — built with RenderCV + Claude",
        ],
        notes=(
            "In just the last couple of months, I've built all of these — entirely with "
            "Claude. python-snap7 for industrial PLC communication, ableton-claude to "
            "control Ableton Live with AI, spectrageist for real-time audio visualization, "
            "geistwave — an experimental music LLM trained on my personal library, "
            "charcoaloptimizer for 3D packing, and this CV and talk itself. None of these "
            "would have existed without AI assistance."
        ),
    )

    # ── Section Break: Advice ──
    add_section_slide(
        prs,
        "Advice for Young Engineers",
        "What I wish someone had told me",
        notes=(
            "OK, that's my story. Now let me share some advice. These are the things I "
            "wish someone had told me when I was starting out. Some of it I learned the "
            "hard way."
        ),
    )

    # ── Slide 21: Follow Interesting, Not Money ──
    add_content_slide(
        prs,
        "Follow Interesting, Not Money",
        [
            "The best career decisions were about curiosity",
            "Chasing money made me unhappy",
            "Money follows interesting work — not the other way around",
            "I never had a five-year plan",
        ],
        notes=(
            "My journey wasn't planned. I never had a five-year career strategy. I just "
            "followed what looked cool — security, then software engineering, then AI, "
            "then radio astronomy. Each pivot happened because something interesting "
            "appeared. When I chased money during my consulting years, I was miserable. "
            "When I followed curiosity, the money eventually came anyway."
        ),
    )

    # ── Slide 22: Build Real Understanding ──
    add_content_slide(
        prs,
        "Build Real Understanding",
        [
            "Don't let AI do everything for you",
            "You need to understand how software works",
            "AI amplifies skills — it doesn't replace learning",
            "Everyone with AI has the same baseline",
            "Your understanding is your edge",
        ],
        notes=(
            "Here's the nuance about AI: you still need to understand programming. "
            "AI assistants amplify your skills — they don't replace the need to learn. "
            "Get your hands dirty first, then use these tools to accelerate. Everyone "
            "has access to AI now — that's the baseline. What differentiates you is your "
            "understanding of fundamentals. Your taste, your judgment, your ability to "
            "ask the right questions."
        ),
    )

    # ── Slide 23: Contribute to Open Source ──
    add_content_slide(
        prs,
        "Contribute to Open Source",
        [
            "Start small — fix bugs, improve docs",
            "Teaches more than any course",
            "Builds your portfolio and reputation",
            "Connects you with developers worldwide",
            "Start today — pick a project you use and love",
        ],
        notes=(
            "Open source was my great equalizer. You don't need permission to "
            "contribute. Pick a project you use and love. Fix a bug. Improve the docs. "
            "Package something. It teaches you real-world collaboration, code review, "
            "and how actual software is built and maintained. And it's your public "
            "portfolio — way more valuable than a certificate."
        ),
    )

    # ── Slide 24: The Namibia Advantage ──
    add_content_slide(
        prs,
        "The Namibia Advantage",
        [
            "Same timezone as Europe — no midnight meetings",
            "English-speaking — EU companies need you",
            "Lower cost of living than Amsterdam or London",
            "Remote work is real — I do it from here",
            "The network builds itself if you stay curious",
        ],
        notes=(
            "This isn't a talk about Silicon Valley dreams. Let me make the practical "
            "case for building a tech career from Namibia. You're in almost the same "
            "timezone as Europe — no midnight meetings. Everyone speaks English — European "
            "companies are actively looking for English-speaking developers. Cost of "
            "living is lower. Internet is getting faster. Remote work isn't a fantasy — "
            "I'm literally doing it. The network builds itself if you stay curious, "
            "contribute to open source, and keep showing up."
        ),
    )

    # ── Slide 25: But: Work Hard ──
    add_content_slide(
        prs,
        "But: Work Hard",
        [
            "Competitive field — especially now with AI",
            "Don't coast",
            "Stay curious, keep learning",
            "AI tools are the new baseline — not the finish line",
        ],
        notes=(
            "I don't want to sugarcoat it. This is a competitive field, and it's getting "
            "more competitive with AI. Don't coast. Stay curious. Keep learning. AI tools "
            "are the new baseline — everyone has them. What matters is what you build on "
            "top of that baseline. The people who succeed are the ones who keep pushing "
            "when others stop."
        ),
    )

    # ── Slide 26: Imposter Syndrome is Normal ──
    add_content_slide(
        prs,
        "Imposter Syndrome is Normal",
        [
            "Early confidence = ignorance (Dunning-Kruger)",
            "Later humility = understanding how much there is to know",
            "It never goes away — it just changes shape",
            "If you feel like a fraud sometimes: welcome to the club",
        ],
        notes=(
            "When I was young, I thought I was the smartest kid in the room. The more I "
            "learned, the more I realized how little I knew. Working with brilliant minds "
            "at the SKA, at research institutions, at Spotify — imposter syndrome didn't "
            "fade with experience. It got worse. Here's what I've learned: this is normal. "
            "The Dunning-Kruger effect is real. If you feel like a fraud sometimes, "
            "welcome to the club. We all feel it."
        ),
    )

    # ── Slide 27: The Realistic Path ──
    add_content_slide(
        prs,
        "The Realistic Path",
        [
            "No genius required",
            "No expensive bootcamps needed",
            "Persistence + curiosity + willingness to fail publicly",
            "Find the right path — it changes everything",
            "Being a bad student doesn't define you",
        ],
        notes=(
            "Let me wrap up the advice section. You don't need to be a genius. You don't "
            "need expensive bootcamps. What you need is persistence, curiosity, and the "
            "willingness to fail publicly. Finding the right path changes everything — I "
            "was a terrible student until I found computing. Not knowing what you want is "
            "normal. Being a bad student doesn't define you. What defines you is what you "
            "do once you find the thing that lights you up."
        ),
    )

    # ── Slide 28: Thank You / Q&A ──
    add_end_slide(
        prs,
        "Thank You!",
        [
            "Gijs Molenaar",
            "Software Engineer @ Spotify",
            "",
            "github.com/gijzelaerr",
            "gijs@pythonic.nl",
            "",
            "Questions?",
        ],
        notes=(
            "Thank you so much for listening. I'm happy to take any questions. "
            "You can find me on GitHub, or reach out by email. "
            "I'm around for the rest of the conference — come say hi!"
        ),
    )

    output = "talk.pptx"
    prs.save(output)
    print(f"Saved {output} with {len(prs.slides)} slides")


if __name__ == "__main__":
    main()
